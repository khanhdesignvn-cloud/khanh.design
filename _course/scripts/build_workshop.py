#!/usr/bin/env python3
"""Build the founding-cohort workshop deck.

The deck intentionally uses only native shapes and system fonts so it can
be edited and exported without external assets.
"""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

BLACK = "0B0B0B"
IVORY = "F4EEDF"
GOLD = "C59B63"
MUTED = "A69D8D"
INK = "181713"
WHITE = "FFFFFF"
FONT = "Liberation Sans"

SLIDES = [
    {
        "kicker": "WORKSHOP MIỄN PHÍ · 90 PHÚT",
        "title": "Xây trợ lý AI đầu tiên\ncho doanh nghiệp với Claude",
        "body": "Một công việc thật · Một bản nháp · Một checklist · Một người duyệt",
        "dark": True,
        "type": "cover",
    },
    {
        "kicker": "ĐẦU RA ĐỘC LẬP",
        "title": "Bạn mang về gì sau 90 phút?",
        "items": [
            "Bản đặc tả trợ lý AI nhỏ theo 6 ô",
            "Một prompt có cấu trúc để chạy thử",
            "Checklist Đạt / Cần sửa",
            "Ranh giới dữ liệu và bước duyệt rõ ràng",
        ],
        "dark": False,
        "type": "list",
    },
    {
        "kicker": "05—15 PHÚT",
        "title": "AI rời rạc thường gãy ở 3 điểm",
        "items": [
            "01  Thiếu bối cảnh — AI không biết điều gì là nguồn chính",
            "02  Thiếu tiêu chuẩn — “hay hơn” không phải một tiêu chí",
            "03  Thiếu điểm dừng — bản nháp bị dùng như quyết định",
        ],
        "dark": True,
        "type": "cards",
    },
    {
        "kicker": "TỪ CHAT ĐẾN HỆ THỐNG",
        "title": "Một luồng có kiểm soát",
        "flow": ["Mục tiêu", "Nguồn sạch", "Hướng dẫn", "Bản nháp", "Kiểm tra", "Người duyệt"],
        "body": "Claude hỗ trợ xử lý. Con người giữ trách nhiệm.",
        "dark": False,
        "type": "flow",
    },
    {
        "kicker": "MÔ HÌNH 6 THÀNH PHẦN",
        "title": "Đừng bắt đầu bằng prompt",
        "items": ["Mục tiêu", "Bối cảnh", "Dữ liệu vào", "Đầu ra", "Tiêu chí", "Ranh giới + duyệt"],
        "dark": True,
        "type": "grid",
    },
    {
        "kicker": "CHỌN VIỆC ĐẦU TIÊN",
        "title": "Nhỏ, rõ, kiểm tra được",
        "items": [
            "Lặp lại và có đầu vào rõ",
            "Đầu ra dễ đối chiếu",
            "Rủi ro thấp, có thể dừng",
            "Tránh quyết định pháp lý, tài chính, nhân sự",
        ],
        "dark": False,
        "type": "list",
    },
    {
        "kicker": "DEMO · DỮ LIỆU GIẢ LẬP",
        "title": "Tiệm bánh Mộc",
        "items": [
            "Bánh thủ công theo đơn",
            "Giọng ấm áp, cụ thể, không khoa trương",
            "Không nhận đơn gấp dưới 24 giờ",
            "Giá luôn để [GIÁ] nếu nguồn chưa cung cấp",
        ],
        "body": "Không phải case study kết quả thật.",
        "dark": True,
        "type": "profile",
    },
    {
        "kicker": "DEMO · BƯỚC 1",
        "title": "Tạo không gian bối cảnh",
        "items": [
            "Nguồn nào được ưu tiên?",
            "Điều gì tuyệt đối không được đoán?",
            "Khi thiếu dữ liệu, phải hỏi lại thế nào?",
        ],
        "dark": False,
        "type": "list",
    },
    {
        "kicker": "DEMO · BƯỚC 2",
        "title": "Giao một nhiệm vụ có thể chấm",
        "quote": "Soạn bản nháp trả lời khách hỏi bánh sinh nhật. Chỉ dùng dữ liệu đã cho; hỏi lại phần thiếu; không tự điền giá; kết thúc bằng một bước tiếp theo.",
        "dark": True,
        "type": "quote",
    },
    {
        "kicker": "DEMO · BƯỚC 3",
        "title": "Sửa theo tiêu chí, không sửa theo cảm giác",
        "items": [
            "Đúng dữ liệu nguồn?",
            "Đúng giọng thương hiệu?",
            "Không bịa giá/chính sách?",
            "Có bước tiếp theo phù hợp?",
        ],
        "dark": False,
        "type": "check",
    },
    {
        "kicker": "47—67 PHÚT",
        "title": "Bài thực hành",
        "body": "Chọn 1 việc: FAQ · bản nháp giới thiệu sản phẩm · tóm tắt ghi chú đã ẩn danh",
        "quote": "Mục tiêu không phải tạo câu trả lời hoàn hảo. Mục tiêu là biết vì sao đầu ra đạt hoặc cần sửa.",
        "dark": True,
        "type": "exercise",
    },
    {
        "kicker": "WORKSHEET",
        "title": "Điền 6 ô trước khi mở Claude",
        "items": ["1. Mục tiêu", "2. Bối cảnh", "3. Dữ liệu", "4. Đầu ra", "5. Tiêu chí", "6. Ranh giới"],
        "dark": False,
        "type": "grid",
    },
    {
        "kicker": "PROMPT KHUNG",
        "title": "Một cấu trúc có thể tái sử dụng",
        "quote": "Vai trò hỗ trợ… Mục tiêu… Bối cảnh… Dữ liệu… Đầu ra… Tiêu chí… Nếu thiếu thông tin… Không được…",
        "body": "Thay nội dung trong từng phần; không dán credential hoặc dữ liệu khách hàng thật.",
        "dark": True,
        "type": "quote",
    },
    {
        "kicker": "CHẠY THỬ",
        "title": "Đánh dấu nơi cần con người",
        "items": [
            "Câu nào cần đối chiếu nguồn?",
            "Phần nào cần quyết định kinh doanh?",
            "Điều gì phải chuyển cấp?",
            "Chưa gửi bản nháp ra bên ngoài.",
        ],
        "dark": False,
        "type": "list",
    },
    {
        "kicker": "67—77 PHÚT",
        "title": "Checklist đánh giá",
        "items": [
            "Đúng thông tin nguồn",
            "Không đoán dữ kiện",
            "Đúng giọng & định dạng",
            "Không vượt quyền",
            "Đã có người chịu trách nhiệm duyệt",
        ],
        "body": "Kết luận: ĐẠT hoặc CẦN SỬA — kèm một lý do cụ thể.",
        "dark": True,
        "type": "check",
    },
    {
        "kicker": "RANH GIỚI",
        "title": "Khi nào phải dừng AI?",
        "items": [
            "Dữ liệu nhạy cảm chưa ẩn danh",
            "Hậu quả pháp lý / tài chính cao",
            "Không có nguồn để kiểm chứng",
            "Người nhận có thể hiểu nháp là quyết định chính thức",
        ],
        "dark": False,
        "type": "list",
    },
    {
        "kicker": "NẾU MUỐN ĐI TIẾP",
        "title": "AI VẬN HÀNH DOANH NGHIỆP",
        "items": [
            "Tuần 1–2 · Giao việc + trợ lý hiểu doanh nghiệp",
            "Tuần 3–4 · Nội dung + bán hàng",
            "Tuần 5–6 · CSKH/vận hành + workflow có duyệt",
        ],
        "body": "6 tuần · học theo dự án thật · không yêu cầu lập trình",
        "dark": True,
        "type": "timeline",
    },
    {
        "kicker": "COHORT SÁNG LẬP",
        "title": "Một lời mời minh bạch",
        "items": [
            "Tối đa 15 học viên",
            "5 suất sáng lập · 3.900.000đ",
            "10 suất chính thức · 4.900.000đ",
            "Chưa gồm phí Claude / công cụ bên thứ ba",
        ],
        "body": "Nộp hồ sơ để cùng đánh giá mức phù hợp — không phải bước thanh toán.",
        "dark": False,
        "type": "pricing",
    },
    {
        "kicker": "BƯỚC TIẾP THEO",
        "title": "Hoàn thiện một vòng nhỏ",
        "items": [
            "01  Chạy lại với một ví dụ đã ẩn danh",
            "02  Ghi một lỗi cụ thể và sửa theo checklist",
            "03  Nếu cần đồng hành 6 tuần, đọc đề cương và nộp hồ sơ",
        ],
        "body": "Câu hỏi?",
        "dark": True,
        "type": "close",
    },
]


def color(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def set_background(slide, value: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color(value)


def add_text(slide, text: str, x: float, y: float, w: float, h: float, *,
             size: int = 24, value: str = INK, bold: bool = False,
             align=PP_ALIGN.LEFT, font: str = FONT, margin: float = 0.02):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(margin)
    frame.margin_top = frame.margin_bottom = Inches(margin)
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color(value)
    return box


def add_rule(slide, y: float, dark: bool) -> None:
    rule = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.65), Inches(y), Inches(12.05), Inches(0.025))
    rule.fill.solid()
    rule.fill.fore_color.rgb = color(GOLD)
    rule.line.fill.background()
    add_text(slide, "khanh.design", 0.67, 7.10, 2.0, 0.22, size=9, value=MUTED if dark else "6B645B")


def add_header(slide, data: dict) -> None:
    dark = data["dark"]
    foreground = IVORY if dark else INK
    add_text(slide, data["kicker"], 0.7, 0.48, 5.8, 0.3, size=10, value=GOLD, bold=True)
    add_text(slide, data["title"], 0.7, 0.95, 11.7, 1.35, size=30 if "\n" not in data["title"] else 28, value=foreground, bold=True)
    add_rule(slide, 6.88, dark)


def add_bullets(slide, items: list[str], dark: bool, y: float = 2.45) -> None:
    foreground = IVORY if dark else INK
    for index, item in enumerate(items):
        top = y + index * (3.72 / max(len(items), 4))
        marker = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.76), Inches(top + 0.12), Inches(0.16), Inches(0.16))
        marker.fill.solid(); marker.fill.fore_color.rgb = color(GOLD); marker.line.fill.background()
        add_text(slide, item, 1.13, top, 10.8, 0.62, size=19 if len(items) <= 4 else 17, value=foreground, bold=False)


def add_grid(slide, items: list[str], dark: bool) -> None:
    foreground = IVORY if dark else INK
    line = "2B2925" if dark else "D7CCB8"
    for index, item in enumerate(items):
        col, row = index % 3, index // 3
        x, y = 0.72 + col * 4.13, 2.48 + row * 1.68
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.75), Inches(1.25))
        card.fill.solid(); card.fill.fore_color.rgb = color("151515" if dark else "FBF8F0")
        card.line.color.rgb = color(line)
        add_text(slide, f"{index + 1:02}", x + 0.22, y + 0.2, 0.45, 0.28, size=10, value=GOLD, bold=True)
        add_text(slide, item, x + 0.22, y + 0.55, 3.25, 0.42, size=16, value=foreground, bold=True)


def add_flow(slide, items: list[str]) -> None:
    for index, item in enumerate(items):
        x = 0.72 + index * 2.02
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.0), Inches(1.62), Inches(1.0))
        box.fill.solid(); box.fill.fore_color.rgb = color(INK if index == len(items) - 1 else "FBF8F0")
        box.line.color.rgb = color(GOLD)
        add_text(slide, item, x + 0.08, 3.31, 1.46, 0.35, size=12, value=IVORY if index == len(items)-1 else INK, bold=True, align=PP_ALIGN.CENTER)
        if index < len(items) - 1:
            add_text(slide, "→", x + 1.66, 3.30, 0.33, 0.3, size=16, value=GOLD, bold=True, align=PP_ALIGN.CENTER)


def add_quote(slide, quote: str, dark: bool, y: float = 2.55) -> None:
    foreground = IVORY if dark else INK
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.75), Inches(y), Inches(0.08), Inches(2.55))
    line.fill.solid(); line.fill.fore_color.rgb = color(GOLD); line.line.fill.background()
    add_text(slide, f"“{quote}”", 1.18, y + 0.05, 10.7, 2.2, size=22, value=foreground, bold=False)


def add_notes(slide, data: dict) -> None:
    notes = slide.notes_slide.notes_text_frame
    notes.text = (
        f"{data['kicker']} — {data['title'].replace(chr(10), ' ')}\n"
        "Dẫn theo workshop-outline.md. Giữ ví dụ ở trạng thái giả lập; mời người học "
        "đối chiếu nguồn và không chia sẻ dữ liệu nhạy cảm."
    )


def build(output: Path) -> Path:
    deck = Presentation()
    deck.slide_width = Inches(13.333)
    deck.slide_height = Inches(7.5)
    deck.core_properties.title = "Xây trợ lý AI đầu tiên cho doanh nghiệp với Claude"
    deck.core_properties.subject = "Workshop miễn phí — AI VẬN HÀNH DOANH NGHIỆP"
    deck.core_properties.author = "khanh.design"
    deck.core_properties.keywords = "Claude, AI, doanh nghiệp, workshop, trợ lý AI"

    blank = deck.slide_layouts[6]
    for number, data in enumerate(SLIDES, start=1):
        slide = deck.slides.add_slide(blank)
        set_background(slide, BLACK if data["dark"] else IVORY)
        add_header(slide, data)
        dark = data["dark"]
        kind = data["type"]

        if kind == "cover":
            add_text(slide, data["body"], 0.74, 3.35, 9.8, 0.6, size=17, value=MUTED)
            add_text(slide, "90’", 10.55, 3.0, 1.65, 1.0, size=42, value=GOLD, bold=True, align=PP_ALIGN.CENTER)
            add_text(slide, "HỌC + LÀM", 10.55, 4.05, 1.65, 0.32, size=9, value=IVORY, bold=True, align=PP_ALIGN.CENTER)
        elif kind in {"list", "check", "profile", "pricing", "timeline", "close", "cards"}:
            add_bullets(slide, data["items"], dark)
        elif kind == "grid":
            add_grid(slide, data["items"], dark)
        elif kind == "flow":
            add_flow(slide, data["flow"])
        elif kind == "quote":
            add_quote(slide, data["quote"], dark)
        elif kind == "exercise":
            add_text(slide, data["body"], 0.75, 2.50, 11.2, 0.7, size=20, value=IVORY, bold=True)
            add_quote(slide, data["quote"], dark, 3.45)

        if data.get("body") and kind not in {"cover", "exercise", "flow"}:
            add_text(slide, data["body"], 0.75, 6.12, 11.4, 0.48, size=12, value=MUTED if dark else "6B645B")
        if kind == "flow":
            add_text(slide, data["body"], 0.75, 5.22, 11.4, 0.45, size=15, value="6B645B", bold=True)

        add_text(slide, f"{number:02}", 12.0, 7.08, 0.7, 0.22, size=9, value=MUTED if dark else "6B645B", align=PP_ALIGN.RIGHT)
        add_notes(slide, data)

    output.parent.mkdir(parents=True, exist_ok=True)
    deck.save(output)
    return output


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    default = Path(__file__).resolve().parents[1] / "dist" / "04-09-2026-Workshop-Xay-tro-ly-AI-dau-tien-voi-Claude.pptx"
    parser.add_argument("--output", type=Path, default=default)
    args = parser.parse_args()
    path = build(args.output.resolve())
    print(f"Built {path} ({len(SLIDES)} slides)")


if __name__ == "__main__":
    main()
