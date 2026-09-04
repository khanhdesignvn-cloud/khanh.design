import json
import re
import subprocess
import sys
from pathlib import Path

import pytest

ROOT = Path(__file__).resolve().parents[1]
COURSE = ROOT / "course"
WORKSHOP = COURSE / "workshop"
LAUNCH = COURSE / "launch"
ONBOARDING = COURSE / "onboarding"
DIST = COURSE / "dist"

FORBIDDEN = (
    "x10",
    "gấp 10 doanh thu",
    "80% công việc",
    "chỉ còn hôm nay",
    "cơ hội cuối cùng",
)
PROGRAM = "AI VẬN HÀNH DOANH NGHIỆP"


def read(path: Path) -> str:
    return path.read_text(encoding="utf-8")


def test_launch_and_onboarding_source_package_is_complete_and_responsible():
    required = [
        WORKSHOP / "workshop-outline.md",
        LAUNCH / "content-plan.md",
        LAUNCH / "posts.md",
        LAUNCH / "demo-video-scripts.md",
        LAUNCH / "case-study-template.md",
        LAUNCH / "workshop-registration-copy.md",
        LAUNCH / "email-zalo-sequence.md",
        ONBOARDING / "application-form-spec.json",
        ONBOARDING / "intake-checklist.md",
        ONBOARDING / "privacy-notice.md",
    ]
    assert not [str(path.relative_to(ROOT)) for path in required if not path.exists()]

    corpus = "\n".join(read(path) for path in required if path.suffix == ".md")
    folded = corpus.casefold()
    assert not [claim for claim in FORBIDDEN if claim.casefold() in folded]
    assert PROGRAM in corpus

    plan = read(LAUNCH / "content-plan.md")
    assert len(re.findall(r"^## Tuần [123]$", plan, flags=re.MULTILINE)) == 3

    posts = read(LAUNCH / "posts.md")
    assert len(re.findall(r"^## Bài [1-5] —", posts, flags=re.MULTILINE)) == 5
    assert posts.count("**CTA:**") == 5

    videos = read(LAUNCH / "demo-video-scripts.md")
    assert len(re.findall(r"^## Video [1-3] —", videos, flags=re.MULTILINE)) == 3
    assert videos.count("**CTA:**") == 3

    cases = read(LAUNCH / "case-study-template.md")
    assert len(re.findall(r"^## Template [12] —", cases, flags=re.MULTILINE)) == 2
    assert "Sự đồng ý công bố" in cases
    assert "Không có dữ liệu thì không nêu con số" in cases


def test_application_form_spec_collects_only_needed_data_with_explicit_consent():
    spec = json.loads(read(ONBOARDING / "application-form-spec.json"))
    assert spec["title"] == f"Đăng ký {PROGRAM} — Cohort sáng lập"
    assert spec["response_destination"] == "REVIEW_REQUIRED"
    assert spec["public_response_sheet"] is False
    assert spec["collect_email"] is True
    ids = {field["id"] for field in spec["fields"]}
    assert {
        "full_name",
        "phone",
        "email",
        "industry",
        "role",
        "business_problem",
        "attendance_commitment",
        "data_consent",
    } <= ids
    consent = next(field for field in spec["fields"] if field["id"] == "data_consent")
    assert consent["required"] is True
    assert consent["type"] == "checkbox"
    field_prompts = " ".join(
        f"{field['id']} {field['label']}" for field in spec["fields"]
    ).casefold()
    assert "mật khẩu" not in field_prompts
    assert "api key" not in field_prompts


def test_workshop_builder_creates_branded_pptx(tmp_path: Path):
    pytest.importorskip("pptx")
    output = tmp_path / "workshop.pptx"
    command = [
        sys.executable,
        str(COURSE / "scripts" / "build_workshop.py"),
        "--output",
        str(output),
    ]
    result = subprocess.run(command, cwd=ROOT, capture_output=True, text=True)
    assert result.returncode == 0, result.stderr
    assert output.exists() and output.stat().st_size > 30_000

    from pptx import Presentation

    deck = Presentation(output)
    assert len(deck.slides) >= 16
    all_text = "\n".join(
        shape.text
        for slide in deck.slides
        for shape in slide.shapes
        if hasattr(shape, "text_frame")
    )
    normalized_text = " ".join(all_text.split())
    assert "Xây trợ lý AI đầu tiên cho doanh nghiệp với Claude" in normalized_text
    assert "Bài thực hành" in all_text
    assert "Checklist đánh giá" in all_text
    assert PROGRAM in all_text
    assert "x10" not in all_text.casefold()
    assert "80%" not in all_text.casefold()

    backgrounds = {
        slide.background.fill.fore_color.rgb.__str__()
        for slide in deck.slides
        if slide.background.fill.type is not None
    }
    assert "0B0B0B" in backgrounds
    assert "F4EEDF" in backgrounds
